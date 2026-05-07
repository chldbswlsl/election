"""
보고서 → 발표용 PPTX 자동 생성 (v2 - 레이아웃 개선)
======================================================
- 차트 종횡비 정확 적용해 캡션·텍스트 충돌 방지
- 결합 곡절(p̂) 제거 — 한국 폰트 호환 표기
- 글자 크기·여백 모두 검증된 값 사용

실행: python generate_pptx.py
출력: 노희래_통계분석_발표.pptx (바탕화면 동시 복사)
"""

from datetime import datetime
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

ROOT = Path(__file__).parent
CHARTS = ROOT / "charts"
OUT = ROOT / "노희래_통계분석_발표.pptx"
DESKTOP = Path("C:/Users/노희래/Desktop/노희래_통계분석_발표.pptx")

# ---------- 레이아웃 상수 (16:9 widescreen) ----------
SLIDE_W = 33.867
SLIDE_H = 19.05
MARGIN_X = 1.5
TITLE_TOP = 0.5
TITLE_H = 1.6
RULE_TOP = 2.3
CONTENT_TOP = 3.2

# ---------- 색상·폰트 ----------
KO = "맑은 고딕"
NAVY = RGBColor(0x1F, 0x3A, 0x6B)
DEM = RGBColor(0x2E, 0x7D, 0xD7)
PP = RGBColor(0xE0, 0x31, 0x31)
DARK = RGBColor(0x1A, 0x1F, 0x36)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_BG = RGBColor(0xE7, 0xEE, 0xF7)


def kfont(run, size=14, bold=False, color=DARK):
    """모든 글자에 한국 폰트 적용 — 가장 단순하고 안정적인 방법"""
    run.font.name = KO
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, text, left, top, width, height, *,
             size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    kfont(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, items, left, top, width, height, *,
                size=14, color=DARK, after_pt=8):
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        # item 이 (text, indent_level) 튜플이면 indent 적용
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(after_pt)
        if level == 0:
            prefix = "• "
            run_color = color
        else:
            prefix = "    – "
            run_color = GRAY
        run = p.add_run()
        run.text = prefix + text
        kfont(run, size=size, color=run_color)


def add_chart(slide, name, top, max_width=24.0, max_height=13.0):
    """차트 이미지를 종횡비 유지·중앙정렬해 삽입. 캡션과 충돌 방지."""
    img = Image.open(CHARTS / f"{name}.png")
    aspect = img.width / img.height
    # 가로 우선, 안 맞으면 세로 기준으로 축소
    if max_width / aspect <= max_height:
        w, h = max_width, max_width / aspect
    else:
        w, h = max_height * aspect, max_height
    left = (SLIDE_W - w) / 2
    slide.shapes.add_picture(str(CHARTS / f"{name}.png"),
                             Cm(left), Cm(top),
                             width=Cm(w), height=Cm(h))
    return top + h  # 이미지 하단 좌표 반환


def add_caption(slide, text, top):
    add_text(slide, text, MARGIN_X, top, SLIDE_W - 2 * MARGIN_X, 1.0,
             size=12, color=GRAY, align=PP_ALIGN.CENTER)


def add_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def add_section_header(slide, title, num):
    """슬라이드 상단 타이틀 + 구분선"""
    add_text(slide, f"{num}  {title}",
             MARGIN_X, TITLE_TOP, SLIDE_W - 2 * MARGIN_X, TITLE_H,
             size=22, bold=True, color=NAVY)
    # 구분선 (얇은 직사각형)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Cm(MARGIN_X), Cm(RULE_TOP),
                                  Cm(SLIDE_W - 2 * MARGIN_X), Cm(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()


# ============== 슬라이드 빌드 ==============
prs = Presentation()
prs.slide_width = Cm(SLIDE_W)
prs.slide_height = Cm(SLIDE_H)
blank = prs.slide_layouts[6]

# ---------- Slide 1 — 표지 ----------
s = prs.slides.add_slide(blank)
# 상단 색 띠
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Cm(0), Cm(0), Cm(SLIDE_W), Cm(2.0))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
add_text(s, "도제 과제 보고서 / 발표 자료",
         MARGIN_X, 0.4, SLIDE_W - 2 * MARGIN_X, 1.2,
         size=14, color=RGBColor(0xFF, 0xFF, 0xFF))
# 메인 타이틀
add_text(s, "2026 서울시장 선거 당선 예측",
         MARGIN_X, 5.5, SLIDE_W - 2 * MARGIN_X, 2.5,
         size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(s, "여론조사 데이터의 통계적 추론",
         MARGIN_X, 8.3, SLIDE_W - 2 * MARGIN_X, 1.4,
         size=20, color=GRAY, align=PP_ALIGN.CENTER)
# 강조 박스
emp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Cm(8), Cm(11.5), Cm(SLIDE_W - 16), Cm(2.2))
emp.fill.solid(); emp.fill.fore_color.rgb = LIGHT_BG
emp.line.color.rgb = NAVY; emp.line.width = Pt(1.5)
add_text(s, "정원오 당선 예측  ·  99% 신뢰수준  ·  p-value 3.73 × 10⁻¹²",
         MARGIN_X, 12.0, SLIDE_W - 2 * MARGIN_X, 1.2,
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
# 푸터
add_text(s, f"노희래   /   {datetime.now().strftime('%Y년 %m월 %d일')}",
         MARGIN_X, 17.0, SLIDE_W - 2 * MARGIN_X, 1,
         size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_notes(s,
    "안녕하세요 노희래입니다. 오늘 발표 주제는 2026년 6월 3일 서울시장 선거 당선 예측입니다. "
    "PDF 강의자료에서 배운 추론 통계 개념을 실제 여론조사 5건에 적용해서 본인 나름대로 결과를 예측해 봤습니다. "
    "결론부터 미리 말씀드리면 정원오 후보 당선이 99% 신뢰수준에서 통계적으로 유의하게 예측됐습니다.")

# ---------- Slide 2 — 한 줄 결론 ----------
s = prs.slides.add_slide(blank)
add_section_header(s, "한 줄 결론", "01")
# 큰 결론
add_text(s, "정원오 당선 예측",
         MARGIN_X, 5.5, SLIDE_W - 2 * MARGIN_X, 2.5,
         size=44, bold=True, color=DEM, align=PP_ALIGN.CENTER)
add_text(s, "99% 신뢰수준에서 통계적으로 유의 (p-value = 3.73 × 10⁻¹²)",
         MARGIN_X, 8.3, SLIDE_W - 2 * MARGIN_X, 1.4,
         size=20, color=GRAY, align=PP_ALIGN.CENTER)
# 득표율 비교 박스
box_w = 12
box_h = 4
add_text(s, "정원오 (민주당)",
         (SLIDE_W - 2 * box_w - 1) / 2, 11, box_w, 1,
         size=16, bold=True, color=DEM, align=PP_ALIGN.CENTER)
add_text(s, "42.83 %",
         (SLIDE_W - 2 * box_w - 1) / 2, 12, box_w, 2,
         size=36, bold=True, color=DEM, align=PP_ALIGN.CENTER)
add_text(s, "오세훈 (국민의힘)",
         (SLIDE_W - 2 * box_w - 1) / 2 + box_w + 1, 11, box_w, 1,
         size=16, bold=True, color=PP, align=PP_ALIGN.CENTER)
add_text(s, "33.64 %",
         (SLIDE_W - 2 * box_w - 1) / 2 + box_w + 1, 12, box_w, 2,
         size=36, bold=True, color=PP, align=PP_ALIGN.CENTER)
add_text(s, "표본수 n_eff = 4,200  /  격차 +9.19 %p",
         MARGIN_X, 15.5, SLIDE_W - 2 * MARGIN_X, 1.2,
         size=16, color=DARK, align=PP_ALIGN.CENTER)
add_notes(s,
    "결론부터 말씀드리겠습니다. 분석 결과 정원오 후보의 당선이 99% 신뢰수준에서 통계적으로 유의합니다. "
    "p-value 가 3.73 곱하기 10의 마이너스 12승 으로 사실상 0 입니다. "
    "예상 득표율은 정원오 42.83%, 오세훈 33.64% 로 격차 +9.19%p 입니다. "
    "이게 어떻게 나온 결과인지 지금부터 설명드리겠습니다.")

# ---------- Slide 3 — 과제 목표 ----------
s = prs.slides.add_slide(blank)
add_section_header(s, "과제 목표 — 두 가지 핵심 질문", "02")
add_bullets(s, [
    "Q1. 단일 여론조사로 '몇 % 지지'라고 단정할 수 있는가?",
    ("못한다. 표본 비율 p 주변에 '신뢰구간' 만 존재 (PDF p44)", 1),
    ("p ± 1.96 · √(p(1-p)/n)", 1),
    "",
    "Q2. 여러 조사를 합치면 더 좁은 추정이 가능한가?",
    ("가능하다. 중심극한정리에 의해 σ → σ/√n (PDF p28)", 1),
    ("표본 N배 → 신뢰구간 폭 √N배 좁아짐", 1),
    "",
    "이 두 원리를 그대로 코드(analyze.py)로 구현해 서울시장 5건을 분석",
], MARGIN_X + 0.5, 4.0, SLIDE_W - 2 * MARGIN_X - 1, 14, size=16)
add_notes(s,
    "과제에서 풀어야 할 두 핵심 질문입니다. "
    "첫째, 한 번 여론조사로 모집단 진짜 지지율을 단정할 수 있는가. 답은 아니다. "
    "둘째, 여러 조사를 합치면 정확해지는가. 답은 가능하다. "
    "이 두 원리를 파이썬 코드로 구현했습니다.")

# ---------- Slide 4 — 통계 ① 신뢰구간 ----------
s = prs.slides.add_slide(blank)
add_section_header(s, "통계 개념 ① — 모비율 신뢰구간 (PDF p44)", "03")
# 핵심 공식 강조
formula = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Cm(5), Cm(4), Cm(SLIDE_W - 10), Cm(2))
formula.fill.solid(); formula.fill.fore_color.rgb = LIGHT_BG
formula.line.color.rgb = NAVY
add_text(s, "p ± 1.96 · √( p (1-p) / n )",
         MARGIN_X, 4.3, SLIDE_W - 2 * MARGIN_X, 1.4,
         size=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
# 설명 + 손계산
add_bullets(s, [
    "왜 1.96 인가? — 표준정규분포 양쪽 꼬리 합 5%를 자르는 임계값 (PDF p15)",
    "PDF p9: '신뢰는 욕먹지 않을 만큼' — 95% = 100번 중 95번 맞을 표기",
], MARGIN_X + 0.5, 6.5, SLIDE_W - 2 * MARGIN_X - 1, 2.5, size=15)
# 손계산 박스
calc_top = 9.5
calc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Cm(MARGIN_X), Cm(calc_top),
                          Cm(SLIDE_W - 2 * MARGIN_X), Cm(7))
calc.fill.solid(); calc.fill.fore_color.rgb = RGBColor(0xF9, 0xFA, 0xFB)
calc.line.color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
add_text(s, "  손계산 검증 — 5/2 SBS 조사 (n=800, 정원오 41%)",
         MARGIN_X + 0.3, calc_top + 0.3, SLIDE_W - 2 * MARGIN_X - 0.6, 1,
         size=14, bold=True, color=NAVY)
calc_lines = [
    "SE = √(0.41 × 0.59 / 800) = 0.01740",
    "MoE = 1.96 × 0.01740 = 0.0341 = 3.41 %p   ← 언론사 발표 ±3.5%p 와 일치 ✓",
    "95% CI(정원오) = [37.6%, 44.4%]               ← analyze.py 출력과 동일 ✓",
]
for i, line in enumerate(calc_lines):
    add_text(s, "  " + line,
             MARGIN_X + 0.3, calc_top + 1.2 + i * 1.2,
             SLIDE_W - 2 * MARGIN_X - 0.6, 1.1,
             size=14, color=DARK)
add_notes(s,
    "첫 번째 통계 개념은 모비율 신뢰구간입니다. 공식 자체는 PDF p44 그대로입니다. "
    "1.96 이라는 숫자는 표준정규분포에서 양쪽 꼬리 5%를 자를 때 임계값입니다. "
    "5월 2일 SBS 조사 데이터로 손계산해 봤더니, "
    "오차범위가 3.41%p 로 나와서 언론사 발표값 ±3.5%p 와 일치했고, "
    "신뢰구간도 코드 출력과 동일했습니다.")

# ---------- Slide 5 — 통계 ② 가설검정 ----------
s = prs.slides.add_slide(blank)
add_section_header(s, "통계 개념 ② — 가설검정과 p-value (PDF p59)", "04")
# 가설 + 공식
add_text(s, "H₀ : 두 후보 동률      vs      H₁ : 정원오 > 오세훈   (단측 검정)",
         MARGIN_X, 4.0, SLIDE_W - 2 * MARGIN_X, 1.3,
         size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
formula = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Cm(5), Cm(5.7), Cm(SLIDE_W - 10), Cm(2))
formula.fill.solid(); formula.fill.fore_color.rgb = LIGHT_BG
formula.line.color.rgb = NAVY
add_text(s, "Z = (p_정 − p_오) / SE_Δ        p-value = 1 − Φ(Z)",
         MARGIN_X, 6.0, SLIDE_W - 2 * MARGIN_X, 1.4,
         size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_bullets(s, [
    "두 후보 차이의 분산 (다항분포): Var = (p_A + p_B − (p_A − p_B)²) / n",
    "p-value 의 의미 (PDF p60) — 1종 오류가 일어날 확률, 작을수록 H₀ 기각 강함",
    "유의수준 0.05 (95% 신뢰) 미만이면 'H₁ 채택' 으로 판정",
    "",
    "본 분석 Poll-of-polls 결과:",
    ("Z = +6.85   /   p-value = 3.73 × 10⁻¹²", 1),
    ("→ 99% 신뢰수준에서도 H₀(동률) 기각 — 강한 증거", 1),
], MARGIN_X + 0.5, 8.5, SLIDE_W - 2 * MARGIN_X - 1, 9, size=15)
add_notes(s,
    "두 번째 개념은 가설검정입니다. 두 후보 지지율이 통계적으로 정말 다른지 검증합니다. "
    "귀무가설은 두 후보 동률, 대립가설은 정원오 우세입니다. "
    "단측 검정으로 진행했고, 다항분포 공분산을 반영한 SE 공식을 사용했습니다. "
    "본 분석 통합 결과 Z 가 +6.85로 임계값 1.96과 2.58을 모두 한참 넘었고, "
    "p-value 가 사실상 0이라 99% 신뢰수준에서도 H₀를 강하게 기각합니다.")

# ---------- Slide 6 — 데이터 ----------
s = prs.slides.add_slide(blank)
add_section_header(s, "분석 데이터 — 서울시장 5건 여론조사", "05")
table_data = [
    ["조사일", "의뢰처", "조사기관", "n", "정원오", "오세훈"],
    ["2025-12-15", "MBC", "코리아리서치", "800", "34.0%", "36.0%"],
    ["2026-02-10", "MBC", "코리아리서치", "800", "40.0%", "36.0%"],
    ["2026-04-20", "펜앤마이크", "여론조사공정", "1,000", "49.5%", "30.9%"],
    ["2026-04-28", "MBC", "코리아리서치", "800", "48.0%", "32.0%"],
    ["2026-05-02", "SBS", "입소스", "800", "41.0%", "34.0%"],
    ["통합 (Poll-of-polls)", "—", "—", "4,200", "42.83%", "33.64%"],
]
table_top = 4.5
table_h = 11
table = s.shapes.add_table(rows=len(table_data), cols=6,
                           left=Cm(2.5), top=Cm(table_top),
                           width=Cm(SLIDE_W - 5), height=Cm(table_h)).table
col_widths = [4.5, 4.5, 5.5, 3.5, 5, 5]
for ci, w in enumerate(col_widths):
    table.columns[ci].width = Cm(w)
for ri, row in enumerate(table_data):
    for ci, val in enumerate(row):
        cell = table.cell(ri, ci)
        cell.text = ""
        cell.text_frame.margin_left = Cm(0.1)
        cell.text_frame.margin_right = Cm(0.1)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = val
        is_header = ri == 0
        is_total = ri == len(table_data) - 1
        kfont(run, size=14, bold=(is_header or is_total),
              color=RGBColor(0xFF, 0xFF, 0xFF) if is_header else DARK)
        if is_header:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        elif is_total:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_BG
add_text(s,
    "출처: 코리아리서치, 입소스, 여론조사공정 등 (모두 중앙선거여론조사심의위 등록)",
    MARGIN_X, table_top + table_h + 0.3,
    SLIDE_W - 2 * MARGIN_X, 1, size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_notes(s,
    "사용한 데이터 5건입니다. 작년 12월부터 올 5월까지 약 5개월간의 변화를 보여줍니다. "
    "MBC 의뢰가 3건으로 가장 많고, SBS 와 펜앤마이크 의뢰가 각 1건씩입니다. "
    "통합 표본 수는 4,200명, 가중평균 지지율은 정원오 42.83%, 오세훈 33.64% 입니다. "
    "보시면 12월 박빙에서 4월 말부터 정원오 우세로 추세가 명확해집니다.")

# ---------- Slide 7 — 결과 ① 추세 ----------
s = prs.slides.add_slide(blank)
add_section_header(s, "결과 ① — 지지율 추세 + 95% 신뢰구간 밴드", "06")
img_bottom = add_chart(s, "trend", CONTENT_TOP, max_width=24, max_height=12.5)
add_caption(s,
    "12월 박빙(Case B) → 4월 말부터 두 신뢰구간이 분리되어 정원오 우세 확정(Case A)",
    img_bottom + 0.3)
add_notes(s,
    "첫 번째 결과 차트입니다. 시간에 따른 지지율 변화를 보여줍니다. "
    "파란 선이 정원오, 빨간 선이 오세훈 후보, 옅은 색 밴드가 95% 신뢰구간입니다. "
    "12월에는 두 밴드가 겹쳐서 PDF p45의 Case B 박빙 상황이었지만, "
    "4월 말부터는 두 밴드가 완전히 분리되어 Case A 우세 확정으로 전환됩니다. "
    "추세상 정원오의 우위가 점점 굳혀지는 모습입니다.")

# ---------- Slide 8 — 결과 ② Forest plot ----------
s = prs.slides.add_slide(blank)
add_section_header(s, "결과 ② — 조사별 신뢰구간 비교 (Forest Plot)", "07")
img_bottom = add_chart(s, "ci_compare", CONTENT_TOP, max_width=24, max_height=12.5)
add_caption(s,
    "에러바가 겹치면 박빙(Case B), 분리되면 우세 확정(Case A) — PDF p45",
    img_bottom + 0.3)
add_notes(s,
    "두 번째 차트는 조사별 신뢰구간을 가로 막대로 보여주는 Forest Plot 입니다. "
    "각 행이 한 조사이고, 두 후보의 95% 신뢰구간이 서로 겹치는지 분리되는지가 핵심입니다. "
    "위쪽 두 행(2025-12, 2026-02)은 겹쳐있어서 박빙이고, "
    "아래 세 행(4월~5월)은 분리되어 통계적 우세입니다. "
    "오른쪽에 초록색 '오차범위 밖' 으로 표시된 게 우세 확정입니다.")

# ---------- Slide 9 — 결과 ③ CLT ----------
s = prs.slides.add_slide(blank)
add_section_header(s, "결과 ③ — 합치면 좁아진다 (중심극한정리)", "08")
img_bottom = add_chart(s, "clt_effect", CONTENT_TOP, max_width=24, max_height=11)
add_caption(s,
    "단일 조사 MoE ±3.41%p → 5건 합산 ±1.50%p   (2.29배 좁아짐 = √5.25)",
    img_bottom + 0.3)
add_notes(s,
    "세 번째 차트는 PDF p28 중심극한정리의 시각적 의미를 보여줍니다. "
    "왼쪽이 단일 조사 800명일 때, 오른쪽이 5건 합쳐 4,200명일 때입니다. "
    "표본을 5.25배 늘리면 분포 폭이 √5.25, 즉 2.29배 좁아져, "
    "두 후보 분포의 겹치는 영역이 거의 사라집니다. "
    "이게 여러 조사를 통합 분석해야 하는 이유이고, 이번 분석의 핵심 트릭입니다.")

# ---------- Slide 10 — 결과 ④ 가설검정 ----------
s = prs.slides.add_slide(blank)
add_section_header(s, "결과 ④ — 가설검정 시각화", "09")
img_bottom = add_chart(s, "hypothesis", CONTENT_TOP, max_width=24, max_height=11.5)
add_caption(s,
    "관측 Z = +6.85 → 임계값 1.96(α=0.05)·2.58(α=0.01) 모두 한참 초과 → 99% H₀ 기각",
    img_bottom + 0.3)
add_notes(s,
    "마지막 차트는 가설검정의 시각적 표현입니다. "
    "표준정규분포 위에서 관측 Z 가 어디 있는지 보여줍니다. "
    "Z 가 +6.85로 임계값 1.96과 2.58을 모두 한참 넘어, 빨간 꼬리 면적이 사실상 0 입니다. "
    "이게 p-value 가 3.73 × 10⁻¹² 인 시각적 근거입니다. "
    "결론, H₀ 즉 두 후보 동률 가설을 99% 신뢰수준에서도 강하게 기각합니다.")

# ---------- Slide 11 — 한계 ----------
s = prs.slides.add_slide(blank)
add_section_header(s, "한계와 주의사항 — 솔직한 자기 평가", "10")
add_bullets(s, [
    "양자대결 가정 — 후보 등록(5월 중) 후 다자 구도 가능 (조국혁신당, 무소속)",
    "비응답 편향 (PDF p4) — 응답률 10~12%, 무응답자가 응답자와 다를 수 있음",
    "시간 가중치 미적용 — 12월 조사와 5월 조사를 동등 가중",
    "House Effect 미보정 — 의뢰처별 편향 무시 (펜앤마이크 등 보수성향 매체)",
], MARGIN_X + 0.5, 4.0, SLIDE_W - 2 * MARGIN_X - 1, 6, size=15)
# 강조 박스
warn_top = 11
warn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Cm(MARGIN_X), Cm(warn_top),
                          Cm(SLIDE_W - 2 * MARGIN_X), Cm(6))
warn.fill.solid(); warn.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
warn.line.color.rgb = RGBColor(0xF5, 0x9E, 0x0B)
warn.line.width = Pt(1.5)
add_text(s, "⚠ 가장 큰 변수 — 부동층 결집 (PDF p60 2종 오류)",
         MARGIN_X + 0.5, warn_top + 0.4,
         SLIDE_W - 2 * MARGIN_X - 1, 1.2,
         size=18, bold=True, color=RGBColor(0x92, 0x40, 0x0E))
add_text(s,
    "부동층 약 22%가 막판에 어느 쪽으로 결집하느냐가 결과를 흔들 수 있음.\n"
    "비관 시나리오 (부동층 7:3 → 오세훈) 시 격차 9.19%p → 0.4%p 박빙 가능.",
    MARGIN_X + 0.5, warn_top + 1.8,
    SLIDE_W - 2 * MARGIN_X - 1, 3,
    size=14, color=RGBColor(0x78, 0x35, 0x0F))
add_notes(s,
    "한계도 솔직하게 말씀드리겠습니다. 양자대결 가정, 비응답 편향, 시간 가중치 미적용, House Effect 미보정 같은 이슈가 있습니다. "
    "특히 가장 큰 변수는 노란 박스의 부동층 결집입니다. "
    "약 22% 부동층이 막판에 어느 쪽으로 흘러가느냐에 따라 결과가 흔들릴 수 있습니다. "
    "비관 시나리오에서는 9%p 격차가 0.4%p 박빙까지 좁혀집니다. "
    "p-value 가 작다고 결과가 확정되는 건 아니며, 모집단의 마음이 6월 3일까지 안 바뀐다는 보장은 다른 문제임을 인정합니다.")

# ---------- Slide 12 — 검증 + 마무리 ----------
s = prs.slides.add_slide(blank)
add_section_header(s, "6/3 이후 자체 검증 + 정리", "11")
add_text(s, "검증 체크리스트",
         MARGIN_X + 0.5, 4.0, SLIDE_W - 2 * MARGIN_X - 1, 1,
         size=18, bold=True, color=NAVY)
add_bullets(s, [
    "정원오 실제 득표율이 95% CI [41.3%, 44.3%] 안에 들어왔는가?",
    "오세훈 실제 득표율이 95% CI [32.2%, 35.1%] 안에 들어왔는가?",
    "1위 적중인가? — 가장 큰 검증",
    "빗나갔다면 원인은? — 부동층 결집 / 비응답 편향 / 의뢰처 편향",
], MARGIN_X + 0.5, 5.2, SLIDE_W - 2 * MARGIN_X - 1, 6, size=15)
add_text(s, "라이브 사이트 + 코드 저장소",
         MARGIN_X + 0.5, 11.5, SLIDE_W - 2 * MARGIN_X - 1, 1,
         size=18, bold=True, color=NAVY)
add_bullets(s, [
    "사이트:  https://chldbswlsl.github.io/election/",
    "코드:    https://github.com/chldbswlsl/election",
], MARGIN_X + 0.5, 12.7, SLIDE_W - 2 * MARGIN_X - 1, 4, size=14)
add_text(s, "감사합니다.",
         MARGIN_X, 17, SLIDE_W - 2 * MARGIN_X, 1.2,
         size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_notes(s,
    "마지막입니다. 6월 3일 개표 결과가 나오면 다음 항목으로 본 모형을 검증할 계획입니다. "
    "각 후보 실제 득표율이 95% 신뢰구간 안에 들어왔는지, 1위 적중인지가 핵심입니다. "
    "빗나간다면 한계 슬라이드의 어느 항목이 결정적이었는지가 후속 학습 주제입니다. "
    "분석 코드와 인터랙티브 사이트는 깃허브에 공개되어 있습니다. "
    "이상으로 발표를 마치겠습니다. 감사합니다. 질문 있으시면 답변드리겠습니다.")

# 저장
import shutil
try:
    prs.save(OUT)
    print(f"  ✓ {OUT.name} 생성 완료 ({OUT.stat().st_size:,} bytes)")
except PermissionError:
    alt = OUT.with_name(OUT.stem + f"_v{datetime.now().strftime('%H%M%S')}" + OUT.suffix)
    prs.save(alt)
    print(f"  ⚠ 기존 파일 열림 — 대체 이름 저장: {alt.name}")
    OUT_LATEST = alt
else:
    OUT_LATEST = OUT

try:
    shutil.copy(OUT_LATEST, DESKTOP)
    print(f"  ✓ 데스크톱 복사: {DESKTOP}")
except PermissionError:
    print(f"  ⚠ 데스크톱 파일이 PowerPoint 에 열려있어 복사 실패. 닫고 다시 실행하세요.")
print(f"  슬라이드 {len(prs.slides)}장 / 차트 4개 / 발표 노트 12개")
